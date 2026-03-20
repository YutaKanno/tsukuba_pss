"""投手分析 PPTX 生成"""
import io

import matplotlib
matplotlib.use( 'Agg' )
import matplotlib.pyplot as plt
import pandas as pd

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from PIL import Image as PILImage

from pitching.calc_ptList import calc_ptList
from pitching.calc_stats import (
    calc_stats, convert_stats_dict_to_df, calc_overallStats, calc_appearance_history,
)
from pitching.plot_statsTable       import plot_statsTable
from pitching.plot_overallStatsTable import plot_overallStatsTable
from pitching.plot_pt_pieChart      import pt_pieChart
from pitching.plot_courceDist       import course_distPlot
from pitching.plot_courseDetail     import course_detailPlot
from pitching.plot_battedBall       import batted_ball_plot
from pitching.plot_velocityDist     import velocity_dist_plot, batted_type_plot
from pitching.plot_appearanceHistory import plot_appearance_history


# ── レイアウト定数 ─────────────────────────────────────────────────
SLIDE_W   = Inches( 8.27 )   # A4 portrait
SLIDE_H   = Inches( 11.69 )
MARGIN    = Inches( 0.55 )
CONTENT_W = SLIDE_W - 2 * MARGIN
GAP       = Inches( 0.15 )
SECTION_H = Inches( 0.28 )

COLOR_HEADER  = RGBColor( 0x1A, 0x3A, 0x5C )
COLOR_SECTION = RGBColor( 0x2C, 0x3E, 0x50 )
COLOR_LABEL   = RGBColor( 0x1A, 0x3A, 0x5C )
COLOR_WHITE   = RGBColor( 0xFF, 0xFF, 0xFF )
COLOR_SUB     = RGBColor( 0xB0, 0xC4, 0xD8 )


# ── ユーティリティ ────────────────────────────────────────────────
def _fig_to_buf( fig, dpi: int = 180 ) -> io.BytesIO:
    buf = io.BytesIO()
    fig.savefig( buf, format = 'png', dpi = dpi, bbox_inches = 'tight' )
    plt.close( fig )
    buf.seek( 0 )
    return buf


def _image_height( buf: io.BytesIO, display_width ) -> Emu:
    PILImage.MAX_IMAGE_PIXELS = None  # decompression bomb チェックを無効化
    buf.seek( 0 )
    iw, ih = PILImage.open( buf ).size
    buf.seek( 0 )
    return Emu( int( display_width * ih / iw ) )


def _add_rect( slide, left, top, width, height, fill_color, no_line = True ):
    shape = slide.shapes.add_shape( 1, left, top, width, height )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if no_line:
        shape.line.fill.background()
    return shape


def _add_text( slide, text, left, top, width, height,
               font_size, bold = False, color = None, word_wrap = False ):
    txb = slide.shapes.add_textbox( left, top, width, height )
    tf  = txb.text_frame
    tf.word_wrap = word_wrap
    p   = tf.paragraphs[ 0 ]
    run = p.add_run()
    run.text            = text
    run.font.size       = Pt( font_size )
    run.font.bold       = bold
    run.font.color.rgb  = color or RGBColor( 0, 0, 0 )


# ── スライドキャンバス ────────────────────────────────────────────
class _Canvas:
    """スライドへの順次配置。オーバーフロー時に新スライドを追加。"""

    def __init__( self, prs: Presentation, blank_layout, top_start, first_slide = None ):
        self._prs       = prs
        self._layout    = blank_layout
        self._top_start = top_start
        if first_slide is not None:
            self.slide = first_slide
            self._y    = top_start
        else:
            self._new_slide()

    def _new_slide( self ):
        self.slide = self._prs.slides.add_slide( self._layout )
        self._y    = self._top_start

    @property
    def _space_left( self ):
        return SLIDE_H - MARGIN - self._y

    def _ensure( self, needed ):
        if needed > self._space_left:
            self._new_slide()

    def add_image( self, buf: io.BytesIO, width = None ):
        if width is None:
            width = CONTENT_W
        h = _image_height( buf, width )
        self._ensure( h )
        buf.seek( 0 )
        self.slide.shapes.add_picture( buf, MARGIN, self._y, width = width, height = h )
        self._y += h + GAP

    def add_section_heading( self, text: str ):
        self._ensure( SECTION_H + Inches( 0.1 ) )
        _add_rect( self.slide, MARGIN, self._y, Inches( 0.04 ), SECTION_H, COLOR_LABEL )
        _add_text( self.slide, text,
                   MARGIN + Inches( 0.1 ), self._y,
                   CONTENT_W - Inches( 0.1 ), SECTION_H,
                   font_size = 10, bold = True, color = COLOR_LABEL )
        self._y += SECTION_H + Inches( 0.08 )

    def add_text_block( self, text: str, font_size: int = 9 ):
        lines   = text.splitlines() or [ '' ]
        line_h  = Pt( font_size * 1.7 ).pt / 72 * 914400
        total_h = Emu( int( line_h * len( lines ) ) )
        self._ensure( total_h )
        txb = self.slide.shapes.add_textbox( MARGIN, self._y, CONTENT_W, total_h )
        tf  = txb.text_frame
        tf.word_wrap = True
        for i, line in enumerate( lines ):
            p   = tf.paragraphs[ 0 ] if i == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text            = line
            run.font.size       = Pt( font_size )
            run.font.color.rgb  = RGBColor( 0x33, 0x33, 0x33 )
        self._y += total_h + GAP

    def spacer( self, h = Inches( 0.12 ) ):
        self._y += h


# ── ヘッダー ──────────────────────────────────────────────────────
def _page_header_height() -> Emu:
    return Inches( 0.8 )


def _add_page_header( slide, pitcher_name, team_name, start_date, end_date ):
    h = _page_header_height()
    _add_rect( slide, Inches( 0 ), Inches( 0 ), SLIDE_W, h, COLOR_HEADER )
    _add_text( slide, '投手分析レポート',
               MARGIN, Inches( 0.06 ), CONTENT_W, Inches( 0.38 ),
               font_size = 16, bold = True, color = COLOR_WHITE )
    _add_text( slide, f'{pitcher_name}　{team_name}　{start_date} 〜 {end_date}',
               MARGIN, Inches( 0.46 ), CONTENT_W, Inches( 0.28 ),
               font_size = 9, color = COLOR_SUB )


def _add_side_header( slide, label ):
    h = Inches( 0.45 )
    _add_rect( slide, Inches( 0 ), Inches( 0 ), SLIDE_W, h, COLOR_SECTION )
    _add_text( slide, label,
               MARGIN, Inches( 0.10 ), CONTENT_W, Inches( 0.28 ),
               font_size = 12, bold = True, color = COLOR_WHITE )


# ── スライド組み立て ──────────────────────────────────────────────
def _build_summary_slide(
    prs, layout,
    df_p, df_all,
    selected_pitcher, selected_team,
    start_date, end_date,
    pitch_type_colors, comment,
):
    first_slide = prs.slides.add_slide( layout )
    _add_page_header( first_slide, selected_pitcher, selected_team,
                      str( start_date ), str( end_date ) )

    canvas = _Canvas( prs, layout,
                      top_start  = _page_header_height() + Inches( 0.2 ),
                      first_slide = first_slide )

    # 総合スタッツ
    pitcher_label  = f'{selected_pitcher} 全体'
    overall_df = pd.DataFrame( [
        calc_overallStats( df_p,  None, pitcher_label ),
        calc_overallStats( df_p,  '右', 'vs右打者' ),
        calc_overallStats( df_p,  '左', 'vs左打者' ),
        calc_overallStats( df_all, None, '平均値'   ),
    ] )
    _blank_for_avg  = [ '投球数', '打席数', '安打数', '本塁打数', '奪三振数', '四死球数' ]
    _game_only_cols = [ '登板数', '投球回', '失点数', '失点率' ]
    overall_df.loc[ overall_df[ 'index' ] == '平均値',      _blank_for_avg  ] = '--'
    overall_df.loc[ overall_df[ 'index' ] != pitcher_label, _game_only_cols ] = '--'

    canvas.add_section_heading( '総合スタッツ' )
    canvas.add_image( _fig_to_buf( plot_overallStatsTable( overall_df ), dpi = 200 ) )
    canvas.spacer()

    # 球速分布
    buf_vel = velocity_dist_plot( df_p, pitch_type_colors )
    if buf_vel:
        canvas.add_section_heading( '球速分布' )
        canvas.add_image( buf_vel )
        canvas.spacer()

    # 被打球性質
    buf_bat = batted_type_plot( df_p, df_all = df_all, pitcher_name = selected_pitcher )
    if buf_bat:
        canvas.add_section_heading( '被打球性質' )
        canvas.add_image( buf_bat )
        canvas.spacer()

    # 登板履歴
    history_df = calc_appearance_history( df_p )
    buf_hist   = plot_appearance_history( history_df )
    if buf_hist:
        canvas.add_section_heading( '登板履歴' )
        canvas.add_image( buf_hist )
        canvas.spacer()

    # コメント
    if comment and comment.strip():
        canvas.add_section_heading( 'コメント' )
        canvas.add_text_block( comment )


def _build_side_slide( prs, layout, df_p, batter_side, side_label,
                       selected_pitcher, selected_team,
                       start_date, end_date, pitch_type_colors ):
    first_slide = prs.slides.add_slide( layout )
    _add_side_header( first_slide, f'vs {side_label}' )

    canvas = _Canvas( prs, layout,
                      top_start   = Inches( 0.45 ) + Inches( 0.2 ),
                      first_slide = first_slide )

    pt_list = calc_ptList( df_p, batter_side = batter_side )
    if not pt_list:
        canvas.add_text_block( f'vs {side_label} のデータがありません。' )
        return

    # スタッツ表
    stats_dict = {}
    for pt in pt_list:
        stats_dict = calc_stats( df_p, pitch_type = pt,
                                 batter_side = batter_side, stats_dict = stats_dict )
    stats_df  = convert_stats_dict_to_df( stats_dict )
    fig_stats = plot_statsTable( stats_df )
    canvas.add_section_heading( 'スタッツ' )
    canvas.add_image( _fig_to_buf( fig_stats, dpi = 200 ) )
    canvas.spacer()

    # 球種割合（左34%）+ カウント別（右66%）横並び
    _N_B, _N_S = 4, 3
    pie_w      = CONTENT_W * 0.34
    count_w    = CONTENT_W - pie_w
    cell_w     = count_w / _N_B
    header_h   = Inches( 0.18 )

    canvas.add_section_heading( '球種割合 / カウント別球種割合' )

    fig_pie = pt_pieChart( df_p, pitch_type_colors, batter_side = batter_side )
    buf_pie = _fig_to_buf( fig_pie, dpi = 200 )
    pie_h   = _image_height( buf_pie, pie_w )

    # セル高さを実画像から取得（tight保存でfigsize比率とずれるため）
    ref_fig = pt_pieChart( df_p, pitch_type_colors, batter_side = batter_side,
                           S = 0, B = 0, show_labels = False,
                           figsize = ( 3.0, 1.0 ), count_label = '0-0' )
    ref_buf = _fig_to_buf( ref_fig, dpi = 200 )
    cell_h  = _image_height( ref_buf, cell_w )

    total_count_h = header_h + cell_h * _N_S
    combined_h    = max( pie_h, total_count_h )
    canvas._ensure( combined_h )

    # 球種割合（左）
    buf_pie.seek( 0 )
    canvas.slide.shapes.add_picture( buf_pie, MARGIN, canvas._y,
                                     width = pie_w, height = pie_h )

    # カウント別ヘッダー行
    for b in range( _N_B ):
        _add_text( canvas.slide, f'B={b}',
                   MARGIN + pie_w + cell_w * b, canvas._y,
                   cell_w, header_h,
                   font_size = 6, color = RGBColor( 0x55, 0x55, 0x55 ) )

    # カウント別パイチャート（width のみ指定してアスペクト比を保持）
    for s in range( _N_S ):
        for b in range( _N_B ):
            sub_fig = pt_pieChart( df_p, pitch_type_colors,
                                   batter_side = batter_side,
                                   S = s, B = b,
                                   show_labels = False,
                                   figsize     = ( 3.0, 1.0 ),
                                   count_label = f'{b}-{s}' )
            buf_sub = _fig_to_buf( sub_fig, dpi = 200 )
            buf_sub.seek( 0 )
            canvas.slide.shapes.add_picture(
                buf_sub,
                MARGIN + pie_w + cell_w * b,
                canvas._y + header_h + cell_h * s,
                width = cell_w,   # height は省略してアスペクト比自動計算
            )

    canvas._y += combined_h + GAP
    canvas.spacer()

    # コース分布グリッド
    canvas.add_section_heading( 'コース分布 / 詳細 / 打球方向' )
    N      = 5
    pts    = ( pt_list[ :N ] + [ '' ] * N )[ :N ]
    cell_w_course = CONTENT_W / N

    for label_or_fn in ( None, course_distPlot, course_detailPlot, batted_ball_plot ):
        row_bufs = []
        for pt in pts:
            if label_or_fn is None:
                row_bufs.append( pt )
            elif pt:
                buf = label_or_fn( df_p, pitch_type = pt, batter_side = batter_side )
                row_bufs.append( buf )
            else:
                row_bufs.append( None )

        if label_or_fn is None:
            # ラベル行
            label_h = Inches( 0.22 )
            canvas._ensure( label_h )
            for i, pt_name in enumerate( row_bufs ):
                if pt_name:
                    _add_text( canvas.slide, pt_name,
                               MARGIN + cell_w_course * i, canvas._y,
                               cell_w_course, label_h,
                               font_size = 7.5, color = RGBColor( 0x33, 0x33, 0x33 ) )
            canvas._y += label_h
        else:
            max_h = Emu( 0 )
            for buf in row_bufs:
                if buf:
                    max_h = max( max_h, _image_height( buf, cell_w_course ) )
            if max_h == Emu( 0 ):
                continue
            canvas._ensure( max_h )
            for i, buf in enumerate( row_bufs ):
                if buf:
                    h = _image_height( buf, cell_w_course )
                    buf.seek( 0 )
                    canvas.slide.shapes.add_picture(
                        buf,
                        MARGIN + cell_w_course * i, canvas._y,
                        width = cell_w_course, height = h,
                    )
            canvas._y += max_h + Inches( 0.05 )


# ── メイン関数 ────────────────────────────────────────────────────
def generate_pitcher_pptx(
    df_p:             pd.DataFrame,
    df_all:           pd.DataFrame,
    selected_pitcher: str,
    selected_team:    str,
    start_date,
    end_date,
    pitch_type_colors: dict,
    comment:          str = '',
) -> io.BytesIO:

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[ 6 ]  # 完全空白

    # Page 1: 全体サマリー
    _build_summary_slide(
        prs, blank_layout,
        df_p, df_all,
        selected_pitcher, selected_team,
        start_date, end_date,
        pitch_type_colors, comment,
    )

    # Page 2: vs右打者
    _build_side_slide(
        prs, blank_layout, df_p, '右', '右打者',
        selected_pitcher, selected_team,
        start_date, end_date, pitch_type_colors,
    )

    # Page 3: vs左打者
    _build_side_slide(
        prs, blank_layout, df_p, '左', '左打者',
        selected_pitcher, selected_team,
        start_date, end_date, pitch_type_colors,
    )

    buf_out = io.BytesIO()
    prs.save( buf_out )
    buf_out.seek( 0 )
    return buf_out
