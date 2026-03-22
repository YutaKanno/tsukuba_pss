"""全出力モード：スタッツ一覧 + チーム作戦分析 + プレイヤー分析 を1つのPDF/PPTXへ"""
import io
import pandas as pd


# ── PDF マージ ──────────────────────────────────────────────────────────────

def _merge_pdfs( pdf_bufs: list ) -> io.BytesIO:
    """複数の PDF BytesIO を1ファイルにマージして返す。"""
    from pypdf import PdfWriter, PdfReader
    writer = PdfWriter()
    for buf in pdf_bufs:
        if buf is None:
            continue
        buf.seek( 0 )
        reader = PdfReader( buf )
        for page in reader.pages:
            writer.add_page( page )
    out = io.BytesIO()
    writer.write( out )
    out.seek( 0 )
    return out


# ── PDF ページ → PNG 変換（PPTX 用） ──────────────────────────────────────────

def _pdf_to_png_pages( pdf_buf: io.BytesIO, scale: float = 2.0 ) -> list:
    """PDF の各ページを PNG BytesIO のリストに変換する。PyMuPDF (fitz) が必要。"""
    import fitz
    pdf_buf.seek( 0 )
    doc = fitz.open( stream=pdf_buf.read(), filetype='pdf' )
    pages = []
    mat = fitz.Matrix( scale, scale )
    for page in doc:
        pix = page.get_pixmap( matrix=mat, alpha=False )
        buf = io.BytesIO( pix.tobytes( 'png' ) )
        buf.seek( 0 )
        pages.append( buf )
    doc.close()
    return pages


# ── タイトルページ生成 ──────────────────────────────────────────────────────

def _make_title_page( selected_team: str, start_date, end_date ) -> io.BytesIO:
    """横向き A4 のタイトルページ PDF を返す。"""
    import os
    from reportlab.lib.pagesizes import landscape, A4
    from reportlab.lib.units import mm
    from reportlab.lib import colors
    from reportlab.platypus import SimpleDocTemplate, Spacer, Table, TableStyle
    from reportlab.platypus import Paragraph
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    # フォント登録
    font = 'IPAexGothic'
    if font not in pdfmetrics.getRegisteredFontNames():
        fpath = os.path.join( os.path.dirname( __file__ ), 'fonts', 'ipaexg.ttf' )
        if os.path.exists( fpath ):
            try:
                pdfmetrics.registerFont( TTFont( font, fpath ) )
            except Exception:
                font = 'Helvetica'

    PAGE_SIZE = landscape( A4 )
    PAGE_W, PAGE_H = PAGE_SIZE
    MARGIN    = 20 * mm
    C_BG      = colors.HexColor( '#1A3A5C' )
    C_WHITE   = colors.white
    C_SUB     = colors.HexColor( '#B0C4D8' )
    C_ACCENT  = colors.HexColor( '#4A90D9' )

    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=PAGE_SIZE,
        leftMargin=MARGIN, rightMargin=MARGIN,
        topMargin=MARGIN, bottomMargin=MARGIN,
    )

    CONTENT_W = PAGE_W - 2 * MARGIN
    CONTENT_H = PAGE_H - 2 * MARGIN

    def _p( text, size, color=C_WHITE, bold=False, align='CENTER' ):
        al = { 'CENTER': 1, 'LEFT': 0, 'RIGHT': 2 }.get( align, 1 )
        return Paragraph( text, ParagraphStyle(
            'tp', fontName=font, fontSize=size,
            textColor=color, leading=size * 1.6,
            alignment=al, spaceAfter=0,
        ) )

    # ── 外枠テーブル（1セル、ページ全体を覆うデザイン）─────
    accent_bar = Table(
        [ [ '' ] ],
        colWidths=[ CONTENT_W ], rowHeights=[ 2 * mm ],
    )
    accent_bar.setStyle( TableStyle( [
        ( 'BACKGROUND', ( 0, 0 ), ( -1, -1 ), C_ACCENT ),
        ( 'LINEABOVE',  ( 0, 0 ), ( -1, -1 ), 0, C_ACCENT ),
    ] ) )

    spacer_top    = Spacer( 1, CONTENT_H * 0.28 )
    spacer_mid    = Spacer( 1, 10 * mm )
    spacer_small  = Spacer( 1, 6 * mm )
    spacer_bottom = Spacer( 1, CONTENT_H * 0.12 )

    story = [
        spacer_top,
        _p( '打者分析', size=42, bold=True ),
        spacer_mid,
        accent_bar,
        spacer_small,
        _p( selected_team, size=22 ),
        spacer_small,
        _p( f'{ start_date }  〜  { end_date }', size=14, color=C_SUB ),
        spacer_bottom,
    ]

    def _draw_bg( canvas, doc ):
        canvas.saveState()
        canvas.setFillColor( C_BG )
        canvas.rect( 0, 0, PAGE_W, PAGE_H, fill=1, stroke=0 )
        canvas.restoreState()

    doc.build( story, onFirstPage=_draw_bg, onLaterPages=_draw_bg )
    buf.seek( 0 )
    return buf


# ── 統合 PDF 生成 ────────────────────────────────────────────────────────────

def generate_combined_pdf(
    df_all_plays:    pd.DataFrame,
    df_team:         pd.DataFrame,
    df_team_all:     pd.DataFrame,
    selected_team:   str,
    start_date,
    end_date,
    selected_batters: list,
    stats_cols:       tuple,
    team_id:          int,
) -> io.BytesIO:
    """
    スタッツ一覧 + チーム作戦分析 + プレイヤー分析 を1つのPDFとして返す。

    必要ライブラリ: pypdf
    """
    from batter_stats_mode   import _generate_stats_pdf
    from batter_analysis_mode import _generate_strategy_pdf, _generate_player_pdf
    from db import batter_comment_repo

    pdfs = []

    # 0. タイトルページ
    pdfs.append( _make_title_page( selected_team, start_date, end_date ) )

    # 1. スタッツ一覧（全体 / 対右投手 / 対左投手 の3ページ）
    pdfs.append(
        _generate_stats_pdf( df_team, selected_team, start_date, end_date )
    )

    # 2. チーム作戦分析（1ページ）
    pdfs.append(
        _generate_strategy_pdf( df_all_plays, df_team, selected_team, start_date, end_date )
    )

    # 3. プレイヤー分析（選択打者 × 3ページ）
    if selected_batters:
        _all_comments = batter_comment_repo.get_all_comments( team_id )
        comments = { b: _all_comments.get( b, '' ) for b in selected_batters }
        pdfs.append(
            _generate_player_pdf(
                selected_batters, df_team, df_team_all,
                selected_team, start_date, end_date, stats_cols,
                comments=comments,
            )
        )

    return _merge_pdfs( pdfs )


# ── 統合 PPTX 生成 ───────────────────────────────────────────────────────────

def generate_combined_pptx(
    df_all_plays:    pd.DataFrame,
    df_team:         pd.DataFrame,
    df_team_all:     pd.DataFrame,
    selected_team:   str,
    start_date,
    end_date,
    selected_batters: list,
    stats_cols:       tuple,
    team_id:          int,
) -> io.BytesIO:
    """
    スタッツ一覧 + チーム作戦分析 + プレイヤー分析 を1つのPPTXとして返す。

    必要ライブラリ: pypdf, PyMuPDF (fitz), python-pptx
    """
    from pptx import Presentation
    from pptx.util import Mm

    # ── 統合 PDF を生成してページ毎に PNG 化 ──────────────
    combined_pdf = generate_combined_pdf(
        df_all_plays, df_team, df_team_all,
        selected_team, start_date, end_date,
        selected_batters, stats_cols, team_id,
    )
    png_pages = _pdf_to_png_pages( combined_pdf )

    # ── PPTX 組み立て（横向きA4、各ページ = 1スライド）────
    prs = Presentation()
    prs.slide_width  = Mm( 297 )
    prs.slide_height = Mm( 210 )
    blank_layout = prs.slide_layouts[ 6 ]   # 完全空白レイアウト

    for png_buf in png_pages:
        slide = prs.slides.add_slide( blank_layout )
        png_buf.seek( 0 )
        slide.shapes.add_picture(
            png_buf,
            0, 0,
            prs.slide_width, prs.slide_height,
        )

    out = io.BytesIO()
    prs.save( out )
    out.seek( 0 )
    return out
